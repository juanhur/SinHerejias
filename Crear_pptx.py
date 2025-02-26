from pptx import Presentation
from pptx.util import Pt
import re
from datetime import datetime
from MongoDB import guardar_presentacion_en_mongo,obtener_ultimo_proyecto


def extract_text_from_pptx(uploaded_file):
    """
    Lee el contenido de un archivo PowerPoint (.pptx) y extrae toda la información.

    Args:
        archivo_pptx: Ruta del archivo .pptx que deseas leer.

    Returns:
        Un string con todo el contenido extraído de las diapositivas.
    """
    try:
        prs = Presentation(uploaded_file)
        contenido = []

        for i, slide in enumerate(prs.slides):
            contenido.append(f"--- Diapositiva {i + 1} ---")
            
            # Extraer texto de cuadros de texto y títulos
            for shape in slide.shapes:
                if shape.has_text_frame:
                    contenido.append(shape.text)
                
                # Extraer texto de tablas
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        fila = [cell.text for cell in row.cells]
                        contenido.append(" | ".join(fila))  # Unir el contenido de cada celda con separadores

        return "\n".join(contenido)

    except Exception as e:
        return f"Error al leer el archivo PowerPoint: {str(e)}"

#

def actualizar_texto_simple(cell, placeholder, new_text):
    """
    Reemplaza el texto directamente en la celda y aplica formato:
    - Cambia todo el texto a tamaño 11.
    - Interpreta '\n\n', '//n' y '/n/n' como saltos de línea.
    - Aplica negrita a las palabras entre '**'.
    """
    cell.text = ""
    text = cell.text.replace("", new_text)
    # Limpiar el texto de la celda existente
    #cell.text = ""
    # Reemplazar //n y /n/n por saltos de línea convencionales (\n\n)
    text = text.replace("\\n", "\n").replace("/n/n", "\n\n")
    # Dividir el texto por los saltos de línea (\n)
    paragraphs = text.split("\n\n")
    # Si hay párrafos, eliminamos el primer salto de línea al principio
    if len(paragraphs) > 0:
        #cell.text_frame.clear()
        # Procesar el resto de los párrafos
        for paragraph_text in paragraphs:
            paragraph = cell.text_frame.add_paragraph()
            process_paragraph(paragraph_text, paragraph)

    # Asegurarse de que el texto queda formateado correctamente
    #cell.text=cell.text.strip()
    cell.text_frame.word_wrap = True

def process_paragraph(paragraph_text, paragraph):
    """ Procesa cada párrafo, aplicando negrita y formato. """
    # Buscar palabras entre ** usando expresión regular
    words = re.split(r'(\*\*.*?\*\*)', paragraph_text)  # Divide el texto por las palabras con negrita

    for word in words:
        run = paragraph.add_run()

        if word.startswith("**") and word.endswith("**"):
            # Si la palabra está entre **, quitar los asteriscos y aplicar negrita
            run.text = word.strip("**")
            run.font.bold = True
        else:
            # De lo contrario, agregar la palabra normalmente
            run.text = word

        # Añadir un espacio después de cada palabra (para evitar que se junten)
        run.text += " "

        # Aplicar tamaño uniforme de fuente
        run.font.size = Pt(11)



def Crear_Diapositiva(nuevos_textos, ruta_salida,usuario,chat_id,area):
    plantilla = r'plantillas/PLANTILLA1.pptx'
    # Cargar la plantilla de PowerPoint
    prs = Presentation(plantilla)
    año,num=obtener_ultimo_proyecto(area)
    codigo_proyecto=str(area)+"-"+str(año)+"-"+str(num)
    # Iterar a través de las diapositivas en la presentación
    for slide in prs.slides:
        # Iterar a través de los elementos de la diapositiva (shapes)
        for shape in slide.shapes:
            # Verificar si el elemento tiene una tabla
            if shape.has_table:
                print(f"Contiene tabla")
                table = shape.table
                # Iterar sobre todas las filas y columnas de la tabla
                for row_idx, row in enumerate(table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        # Iterar sobre los nuevos textos y verificar qué texto reemplazar
                        for i, texto in enumerate(nuevos_textos):
                            # Buscar y reemplazar palabras clave con formato
                            if "NOMBRE_PROYECTO_IA" in cell.text:
                                actualizar_texto_simple(cell, "NOMBRE_PROYECTO_IA", texto["NOMBRE_PROYECTO_IA"])
                            elif "JUSTIFICACION_IA" in cell.text:
                                actualizar_texto_simple(cell, "JUSTIFICACION_IA", texto["JUSTIFICACION_IA"])
                            elif "COD_IA" in cell.text:
                                actualizar_texto_simple(cell, "COD_IA", codigo_proyecto)
                            elif "OBJETIVOS_IA" in cell.text:
                                actualizar_texto_simple(cell, "OBJETIVOS_IA", texto["OBJETIVOS_IA"])
                            elif "RIESGOS_IA" in cell.text:
                                actualizar_texto_simple(cell, "RIESGOS_IA", texto["RIESGOS_IA"])
                            elif "CLUSTER_IA" in cell.text:
                                actualizar_texto_simple(cell, "CLUSTER_IA", texto["CLUSTER_IA"])
                            elif "PRODUCT_OWNER_IA" in cell.text:
                                actualizar_texto_simple(cell, "PRODUCT_OWNER_IA", texto["PRODUCT_OWNER_IA"])
                            elif "ESPECIFICACIONES_IA" in cell.text:
                                actualizar_texto_simple(cell, "ESPECIFICACIONES_IA", texto["ESPECIFICACIONES_IA"])
                            elif "RECURSOS_IA" in cell.text:
                                actualizar_texto_simple(cell, "RECURSOS_IA", texto["RECURSOS_IA"])
                            elif "CRONOGRAMA_IA" in cell.text:
                                actualizar_texto_simple(cell, "CRONOGRAMA_IA", texto["CRONOGRAMA_IA"])
                            elif "FECHA_IA" in cell.text:
                                actualizar_texto_simple(cell, "FECHA_IA", datetime.now().strftime("%d/%m/%Y"))
    
    # Guardar la presentación con los textos actualizados
    prs.save(ruta_salida)
    guardar_presentacion_en_mongo(ruta_salida, usuario,chat_id)
    return ruta_salida
    
